"""Real offline extraction gate. No database/model clients or source text in JSON.

Run with the prepared backend venv: python scripts/check_extraction_runtime.py
--matrix --output /tmp/extraction-matrix.json. Unavailable routes FAIL, never skip.
"""

from __future__ import annotations

import argparse
import asyncio
import hashlib
import json
import platform
import shutil
import tempfile
from email.message import EmailMessage
from pathlib import Path

from actweave_knowledge.contracts import KnowledgeError
from actweave_knowledge.extraction.contracts import ExtractionLimits, ExtractSetting, ParseProfile
from actweave_knowledge.extraction.registry import default_registry
from actweave_knowledge.extraction.runtime import run_extraction
from actweave_knowledge.extraction.runtime_resources import runtime_digest, runtime_manifest


def matrix_routes():
    return [(mode, extension, item) for item in default_registry().registrations for mode in item.etl_types for extension in item.extensions]


def _fixtures(root: Path) -> dict[str, tuple[Path, str, str, int, set[str]]]:
    from docx import Document
    from ebooklib import epub
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as SheetImage
    from PIL import Image
    from pptx import Presentation

    root.mkdir()
    result = {}
    for extension in (".txt", ".md", ".markdown", ".mdx", ".csv", ".html", ".htm", ".xml"):
        path = root / ("sample" + extension)
        content = "Local matrix text"
        location, warnings = "line", set()
        if extension in (".md", ".markdown", ".mdx"):
            content = "# Local heading\n\nLocal matrix text\n\n![remote](https://example.invalid/image.png)"
            warnings = {"EXTERNAL_IMAGE_NOT_FETCHED"}
        elif extension == ".csv":
            content = "Name,Value\nLocal matrix text,00123\n"
            location, warnings = "row", {"HEADER_INFERRED"}
        elif extension in (".html", ".htm"):
            content = "<html><body><h1>Local heading</h1><p>Local matrix text</p></body></html>"
            location = "element"
        elif extension == ".xml":
            content = '<?xml version="1.0"?><root><p>Local matrix text</p></root>'
            location = "element"
        path.write_text(content, encoding="utf-8")
        result[extension] = (path, "Local matrix text", location, 0, warnings)
    png = root / "red.png"
    Image.new("RGB", (3, 2), "red").save(png)
    word = Document()
    word.add_paragraph("Local matrix text")
    word.add_picture(str(png))
    path = root / "sample.docx"
    word.save(path)
    result[".docx"] = (path, "Local matrix text", "paragraph", 1, set())
    workbook = Workbook()
    workbook.active.append(["Name", "Value"])
    workbook.active.append(["Local matrix text", "00123"])
    workbook.active.add_image(SheetImage(png), "C2")
    path = root / "sample.xlsx"
    workbook.save(path)
    result[".xlsx"] = (path, "Local matrix text", "sheet", 1, {"HEADER_INFERRED"})
    path = root / "sample.xls"
    shutil.copyfile(Path(__file__).parents[1] / "tests/knowledge/fixtures/xlrd/ragged.xls", path)
    result[".xls"] = (path, "l", "sheet", 0, {"HEADER_INFERRED"})
    # Minimal real PDF with a text page; no external writer/executable dependency.
    stream = b"BT /F1 12 Tf 72 720 Td (Local matrix text) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    payload = bytearray(b"%PDF-1.4\n")
    offsets = []
    for index, body in enumerate(objects, 1):
        offsets.append(len(payload))
        payload.extend(f"{index} 0 obj\n".encode() + body + b"\nendobj\n")
    xref = len(payload)
    payload.extend(b"xref\n0 6\n0000000000 65535 f \n")
    for offset in offsets:
        payload.extend(f"{offset:010d} 00000 n \n".encode())
    payload.extend(f"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode())
    path = root / "sample.pdf"
    path.write_bytes(payload)
    result[".pdf"] = (path, "Local matrix text", "page", 0, set())
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Local matrix text"
    slide.placeholders[1].text = "This is a complete local presentation sentence."
    path = root / "sample.pptx"
    deck.save(path)
    result[".pptx"] = (path, "Local matrix text", "slide", 0, set())
    book = epub.EpubBook()
    book.set_identifier("offline-matrix")
    book.set_title("Local matrix text")
    book.set_language("en")
    chapter = epub.EpubHtml(title="Chapter", file_name="chapter.xhtml", lang="en")
    chapter.content = "<h1>Local matrix text</h1><p>This is a complete local book sentence.</p>"
    book.add_item(chapter)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    path = root / "sample.epub"
    epub.write_epub(path, book)
    result[".epub"] = (path, "Local matrix text", "element", 0, set())
    mail = EmailMessage()
    mail["From"] = "sender@example.invalid"
    mail["To"] = "recipient@example.invalid"
    mail["Subject"] = "Local fixture"
    mail.set_content("Local matrix text", cte="base64")
    path = root / "sample.eml"
    path.write_bytes(bytes(mail))
    result[".eml"] = (path, "Local matrix text", "element", 0, set())
    path = root / "sample.msg"
    shutil.copyfile(Path(__file__).parents[1] / "tests/knowledge/fixtures/python-oxmsg/no-attachments.msg", path)
    result[".msg"] = (path, "This is a message", "element", 0, set())
    return result


async def run_matrix(root: Path) -> dict:
    fixtures = await asyncio.to_thread(_fixtures, root / "fixtures")
    rows = []
    for index, (mode, extension, registration) in enumerate(matrix_routes()):
        path, marker, location, expected_assets, expected_warnings = fixtures[extension]
        row = {"mode": mode, "extension": extension, "parser_version": registration.extractor_version, "result": "failed"}
        work_dir = root / f"case-{index}"
        accepted = {}

        async def guard():
            return None

        async def on_asset(asset):
            data = await asyncio.to_thread((work_dir / asset.relative_path).read_bytes)
            assert hashlib.sha256(data).hexdigest() == asset.attachment.ref
            accepted[asset.attachment.ref] = len(data)

        try:
            profile = ParseProfile(etl_type=mode, extractor_id=registration.extractor_id, extractor_version=registration.extractor_version, normalization_version="md-v1", image_policy_version="raster-v1")
            result = await run_extraction(ExtractSetting(source_path=path, original_name=path.name, profile=profile), work_dir=work_dir, limits=ExtractionLimits(), timeout_seconds=120, on_asset=on_asset, guard=guard)
            text = "\n".join(document.page_content for document in result.documents)
            spans = [span for document in result.documents for span in document.source_spans]
            warnings = {warning.code for document in result.documents for warning in document.warnings} | {warning.code for warning in result.warnings}
            assert marker in text and spans and any(location in span.location for span in spans)
            assert len(result.attachments) == len(accepted) == expected_assets
            assert set(accepted) == {asset.ref for asset in result.attachments}
            assert expected_warnings <= warnings
            assert result.source_sha256 == hashlib.sha256(await asyncio.to_thread(path.read_bytes)).hexdigest()
            row.update(result="passed", documents=len(result.documents), source_spans=len(spans), attachments=len(result.attachments), warnings=sorted(warnings))
        except KnowledgeError as error:
            row["reason_code"] = getattr(error, "reason_code", error.code)
        except AssertionError:
            row["reason_code"] = "MATRIX_EXPECTATION_FAILED"
        except Exception:
            row["reason_code"] = "MATRIX_RUNTIME_FAILED"
        rows.append(row)
    return {
        "platform": platform.system(),
        "python": platform.python_version(),
        "resource_digest": runtime_digest(),
        "resource_manifest": runtime_manifest(),
        "matrix": rows,
        "counts": {"passed": sum(row["result"] == "passed" for row in rows), "failed": sum(row["result"] != "passed" for row in rows), "skipped": 0},
    }


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--matrix", action="store_true", required=True)
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()
    with tempfile.TemporaryDirectory(prefix="extraction-matrix-") as directory:
        report = asyncio.run(run_matrix(Path(directory)))
    args.output.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report["counts"]))
    raise SystemExit(bool(report["counts"]["failed"]))


if __name__ == "__main__":
    main()
