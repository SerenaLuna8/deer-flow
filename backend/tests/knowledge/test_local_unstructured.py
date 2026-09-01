"""Real local parsers and stable source attribution; no external services."""

from __future__ import annotations

import hashlib
import importlib
import importlib.util
import json
from email.message import EmailMessage
from pathlib import Path
from types import SimpleNamespace

import pytest
from parsing_test_helpers import make_context, make_parse_profile, make_setting


def elements_converter():
    name = "actweave_knowledge.extraction.unstructured_local.elements"
    assert importlib.util.find_spec("actweave_knowledge.extraction.unstructured_local") is not None, "local adapters not implemented"
    return importlib.import_module(name).elements_to_documents


def test_pptx_missing_page_is_preserved_without_invented_page():
    elements = [SimpleNamespace(text=text, category="Title", metadata=SimpleNamespace(page_number=page)) for text, page in [("第一页", 1), ("无页码但真实存在", None), ("第二页", 2)]]
    docs = elements_converter()(elements, kind="slide")
    assert [d.page_content for d in docs] == ["第一页", "无页码但真实存在", "第二页"]
    assert all("slide" not in s.location for s in docs[1].source_spans)
    assert any(w.code == "SOURCE_POSITION_UNAVAILABLE" for w in docs[1].warnings)
    assert docs[0].source_spans[0].location["slide"] == 1


def test_email_element_is_already_decoded():
    text = "SGVsbG8= 是用户写下的字面值；中文正文。"
    docs = elements_converter()([SimpleNamespace(text=text, category="NarrativeText", metadata=SimpleNamespace())], kind="mail")
    assert docs[0].page_content == text


def test_real_table_html_uses_shared_converter_and_preserves_element_location():
    from unstructured.documents.elements import ElementMetadata, Table, Title

    docs = elements_converter()([Title("Chapter"), Table("Key Value x 00123", metadata=ElementMetadata(page_number=3, text_as_html="<table><tr><th>Key</th><th>Value</th></tr><tr><td>x</td><td>00123</td></tr></table>"))], kind="slide")
    assert "| Key | Value |" in docs[1].page_content
    assert "| x | 00123 |" in docs[1].page_content
    assert docs[1].heading_path == ("Chapter",)
    assert all(s.location["slide"] == 3 and s.location["element"] == 2 for s in docs[1].source_spans)


def test_text_only_table_warns_without_inventing_columns():
    from unstructured.documents.elements import Table

    docs = elements_converter()([Table("x 00123")], kind="mail")
    assert docs[0].page_content == "x 00123"
    assert any(w.code == "TABLE_STRUCTURE_UNAVAILABLE" for w in docs[0].warnings)


def parse(path, tmp_path, etl="unstructured_local"):
    from actweave_knowledge.extraction.processor import ExtractProcessor

    profile = make_parse_profile(path.suffix, etl_type=etl)
    return ExtractProcessor().extract(make_setting(path, profile=profile), make_context(tmp_path / "work"))


@pytest.mark.parametrize("etl", ["dify", "unstructured_local"])
def test_real_pptx(tmp_path, etl):
    from pptx import Presentation

    deck = Presentation()
    for title in ["First slide", "Second slide"]:
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "This is a complete sentence in the presentation."
    path = tmp_path / "slides.pptx"
    deck.save(path)
    docs = parse(path, tmp_path, etl)
    assert "First slide" in "\n".join(d.page_content for d in docs)
    assert {s.location.get("slide") for d in docs for s in d.source_spans} == {1, 2}


@pytest.mark.parametrize("etl", ["dify", "unstructured_local"])
def test_real_epub(tmp_path, etl):
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("fixture")
    book.set_title("Local book")
    book.set_language("en")
    chapter = epub.EpubHtml(title="Chapter One", file_name="ch1.xhtml", lang="en")
    chapter.content = "<h1>Chapter One</h1><p>This is a real local EPUB paragraph.</p>"
    book.add_item(chapter)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    path = tmp_path / "book.epub"
    epub.write_epub(path, book)
    docs = parse(path, tmp_path, etl)
    assert "real local EPUB paragraph" in "\n".join(d.page_content for d in docs)
    assert all("chapter" not in s.location for d in docs for s in d.source_spans)


@pytest.mark.parametrize(
    "text",
    [
        "# C#\n父说明\n## 子节\nList<int> Map<K,V> <IP>\n```cpp\nvector<int> x;\n```\n",
        "## C# ###\n<Component value={dangerous()}/>\n{process.exit(1)}\n`List<int> ![x](https://x)`\n",
        "# 主\n   ~~~~python\n## literal\n![literal](https://no.invalid/x)\n   ~~~~~\n",
        '{"![diagram](https://example.invalid/x)"}\n',
        '<Example value={"![diagram](https://example.invalid/x)"} />\n',
        "<pre>\n![literal](https://example.invalid/x)\n</pre>\n",
    ],
)
def test_real_markdown_literals_and_original_lines(tmp_path, text):
    path = tmp_path / "sample.mdx"
    path.write_text(text)
    docs = parse(path, tmp_path)
    assert "".join(d.page_content for d in docs) == text
    assert all("line" in s.location for d in docs for s in d.source_spans)
    assert not any("ACTWEAVE" in d.page_content for d in docs)


def test_real_eml_mime_decoded_once(tmp_path):
    message = EmailMessage()
    message["From"] = "sender@example.invalid"
    message["To"] = "receiver@example.invalid"
    message["Subject"] = "Local fixture"
    message.set_content("SGVsbG8= 是用户写下的字面值；中文正文。", charset="utf-8", cte="base64")
    path = tmp_path / "message.eml"
    path.write_bytes(bytes(message))
    docs = parse(path, tmp_path)
    assert "SGVsbG8= 是用户写下的字面值；中文正文。" in "\n".join(d.page_content for d in docs)


def test_real_msg_and_pinned_fixture_provenance(tmp_path):
    base = Path(__file__).parent / "fixtures/python-oxmsg"
    provenance = json.loads((base / "provenance.json").read_text())
    for item in provenance["files"]:
        assert hashlib.sha256((base / item["path"]).read_bytes()).hexdigest() == item["sha256"]
    docs = parse(base / "no-attachments.msg", tmp_path)
    assert "This is a message" in "\n".join(d.page_content for d in docs)


def test_real_xml_honors_declared_encoding(tmp_path):
    path = tmp_path / "source.xml"
    path.write_bytes('<?xml version="1.0" encoding="UTF-16"?><root><p>真实 XML 正文</p></root>'.encode("utf-16"))
    docs = parse(path, tmp_path)
    assert "真实 XML 正文" in "\n".join(d.page_content for d in docs)


@pytest.mark.parametrize("encoding", ["utf-8", "utf-16"])
@pytest.mark.parametrize("target", ["file:///etc/passwd", "https://example.invalid/external.dtd"])
def test_xml_dtd_rejected_before_partition(tmp_path, monkeypatch, encoding, target):
    from actweave_knowledge.extraction.contracts import ExtractionError

    path = tmp_path / "hostile.xml"
    path.write_bytes(f'<?xml version="1.0" encoding="{encoding}"?><!DOCTYPE root [<!ENTITY x SYSTEM "{target}">]><root>&x;</root>'.encode(encoding))
    # Adapters themselves must guard XML, independently of processor signatures.
    from actweave_knowledge.extraction.registry import default_registry

    adapter = default_registry().resolve(datasource_type="file", etl_type="unstructured_local", extension=".xml").factory()
    import unstructured.partition.xml

    monkeypatch.setattr(unstructured.partition.xml, "partition_xml", lambda **kw: pytest.fail("unsafe XML reached partition"))
    with pytest.raises(ExtractionError) as error:
        adapter.extract(make_setting(path, profile=make_parse_profile(".xml", etl_type="unstructured_local")), make_context(tmp_path / "work"))
    assert error.value.reason_code == "FORMAT_SIGNATURE_MISMATCH"


@pytest.mark.parametrize("damage", ["drop", "duplicate", "reorder"])
def test_markdown_partition_marker_loss_is_rejected(tmp_path, monkeypatch, damage):
    import unstructured.partition.md
    from actweave_knowledge.extraction.contracts import ExtractionError

    original = unstructured.partition.md.partition_md

    def damaged(**kwargs):
        result = original(**kwargs)
        if damage == "drop":
            return result[1:]
        if damage == "duplicate":
            return result + result[:1]
        return list(reversed(result))

    monkeypatch.setattr(unstructured.partition.md, "partition_md", damaged)
    path = tmp_path / "mapped.md"
    path.write_text("# Heading\n\nOrdinary source text.\n")
    with pytest.raises(ExtractionError) as error:
        parse(path, tmp_path)
    assert error.value.reason_code == "FORMAT_SIGNATURE_MISMATCH"


def test_markdown_external_image_uses_safe_placeholder_and_literal_does_not(tmp_path):
    path = tmp_path / "images.md"
    path.write_text("# Images\n![real](https://example.invalid/r)\n`![literal](https://example.invalid/l)`\n")
    docs = parse(path, tmp_path)
    text = "".join(d.page_content for d in docs)
    assert "（外部图片未获取：real）" in text
    assert "`![literal](https://example.invalid/l)`" in text
    assert any(w.code == "EXTERNAL_IMAGE_NOT_FETCHED" for d in docs for w in d.warnings)


def test_mail_attachments_do_not_enter_auto_partition(tmp_path, monkeypatch):
    import unstructured.partition.auto

    monkeypatch.setattr(unstructured.partition.auto, "partition", lambda **kw: pytest.fail("mail attachment entered auto partition"))
    message = EmailMessage()
    message["From"] = "sender@example.invalid"
    message["To"] = "receiver@example.invalid"
    message["Subject"] = "Attachment fixture"
    message.set_content("This message body must survive without parsing attachments.")
    message.add_attachment(b"unsupported binary file", maintype="application", subtype="octet-stream", filename="input.pdf")
    path = tmp_path / "attachment.eml"
    path.write_bytes(bytes(message))
    docs = parse(path, tmp_path)
    assert "message body must survive" in "\n".join(d.page_content for d in docs)


def test_shared_setting_helper_accepts_local_only_profile(tmp_path):
    profile = make_parse_profile(".eml", etl_type="unstructured_local")
    setting = make_setting(tmp_path / "message.eml", profile=profile)
    assert setting.profile is profile
    assert setting.original_name == "message.eml"


def test_markdown_rejects_literal_marker_reassigned_to_another_source_block(tmp_path, monkeypatch):
    import re

    import unstructured.partition.md
    from actweave_knowledge.extraction.contracts import ExtractionError

    original = unstructured.partition.md.partition_md

    def swapped(**kwargs):
        elements = original(**kwargs)
        found = [(element, match.group()) for element in elements for match in re.finditer(r"ACTWEAVE[A-F0-9]+L\d+Z", element.text)]
        assert len(found) == 2
        (first, first_marker), (second, second_marker) = found
        first.text = first.text.replace(first_marker, second_marker)
        second.text = second.text.replace(second_marker, first_marker)
        return elements

    monkeypatch.setattr(unstructured.partition.md, "partition_md", swapped)
    path = tmp_path / "mapped.md"
    path.write_text("List<int>\n\nMap<K,V>\n")
    with pytest.raises(ExtractionError) as error:
        parse(path, tmp_path)
    assert error.value.reason_code == "FORMAT_SIGNATURE_MISMATCH"
