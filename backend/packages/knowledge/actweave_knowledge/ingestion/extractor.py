"""Synchronous text extraction for the nine supported document formats.

Every function here does blocking file and parser I/O and MUST be called via
``asyncio.to_thread`` from async code. Parser failures, unreadable encodings,
and files without any extractable text all surface as
:class:`KnowledgeError` with ``KNOWLEDGE_PARSE_FAILED``.
"""

from __future__ import annotations

import codecs
import csv
import io
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from ..contracts import KNOWLEDGE_PARSE_FAILED, KNOWLEDGE_QUOTA_EXCEEDED, KnowledgeError

# Zip containers (docx/xlsx) below this declared decompressed total are always
# allowed regardless of the character budget: legitimate media-heavy files
# barely inflate, while decompression bombs need ratios far beyond it.
_ZIP_BYTES_FLOOR = 64 * 1024 * 1024


@dataclass(frozen=True, slots=True)
class ExtractedBlock:
    """One naturally delimited piece of source text with its origin."""

    text: str
    source_position: dict[str, Any] = field(default_factory=dict)


def _parse_failed(message: str) -> KnowledgeError:
    return KnowledgeError(KNOWLEDGE_PARSE_FAILED, message)


@dataclass(slots=True)
class _CharBudget:
    """Cumulative cap on extracted characters.

    ``upload_max_bytes`` does not bound PDF/DOCX/XLSX memory because those
    formats decompress; spending the budget as blocks accumulate aborts a
    bomb long before it can exhaust the worker process.
    """

    limit: int | None
    used: int = 0

    def spend(self, text: str) -> None:
        if self.limit is None:
            return
        self.used += len(text)
        if self.used > self.limit:
            raise KnowledgeError(
                KNOWLEDGE_QUOTA_EXCEEDED,
                f"文档文本超过 {self.limit} 字符（分段配额 × 分段大小），已中止解析",
            )


def _ensure_zip_within_budget(path: Path, budget: _CharBudget, *, kind: str) -> None:
    """Reject a zip container whose declared decompressed size cannot be legit.

    python-docx builds the whole XML DOM and openpyxl loads all shared strings
    up front, so the character budget alone cannot stop a decompression bomb.
    ``zipfile`` never yields more bytes than the central directory declares
    (a mismatch fails the CRC check), so the declared total is a sound bound.
    """

    if budget.limit is None:
        return
    try:
        with zipfile.ZipFile(path) as archive:
            declared_total = sum(info.file_size for info in archive.infolist())
    except (zipfile.BadZipFile, OSError):
        raise _parse_failed(f"{kind} 解析失败") from None
    if declared_total > max(_ZIP_BYTES_FLOOR, 16 * budget.limit):
        raise KnowledgeError(
            KNOWLEDGE_QUOTA_EXCEEDED,
            f"{kind} 解压后大小超过可接受上限，已中止解析",
        )


def extract_blocks(
    path: Path,
    extension: str,
    *,
    max_total_chars: int | None = None,
) -> list[ExtractedBlock]:
    """Extract text blocks from the file at ``path`` routed by ``extension``.

    ``extension`` is the lowercased suffix including the dot (``".pdf"``).
    Raises ``KNOWLEDGE_PARSE_FAILED`` when parsing fails or no block contains
    any non-whitespace text, and ``KNOWLEDGE_QUOTA_EXCEEDED`` as soon as the
    accumulated text passes ``max_total_chars``.
    """

    extractor = _EXTRACTORS.get(extension)
    if extractor is None:
        raise _parse_failed(f"不支持的文件类型 {extension or '(无扩展名)'}")
    blocks = extractor(path, _CharBudget(max_total_chars))
    if not any(block.text.strip() for block in blocks):
        raise _parse_failed("文件没有可提取的文本")
    return blocks


def _decode_text_bytes(data: bytes) -> str:
    """Decode raw text bytes: UTF-16 (BOM), UTF-8, then GB18030."""

    if data.startswith((codecs.BOM_UTF16_LE, codecs.BOM_UTF16_BE)):
        try:
            return data.decode("utf-16")
        except UnicodeDecodeError:
            raise _parse_failed("UTF-16 文本解码失败") from None
    try:
        return data.decode("utf-8-sig")
    except UnicodeDecodeError:
        pass
    try:
        return data.decode("gb18030")
    except UnicodeDecodeError:
        raise _parse_failed("无法识别的文本编码，仅支持 UTF-16、UTF-8 和 GB18030") from None


def _extract_pdf(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            # An empty owner password is common for viewer-only protection.
            if not reader.decrypt(""):
                raise _parse_failed("PDF 已加密，无法解析")
        blocks: list[ExtractedBlock] = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            budget.spend(text)
            blocks.append(ExtractedBlock(text=text, source_position={"page": page_number}))
        return blocks
    except KnowledgeError:
        raise
    except Exception:
        raise _parse_failed("PDF 解析失败") from None


def _extract_docx(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    import docx

    _ensure_zip_within_budget(path, budget, kind="DOCX")
    try:
        document = docx.Document(str(path))
        blocks: list[ExtractedBlock] = []
        for paragraph_number, paragraph in enumerate(document.paragraphs, start=1):
            text = paragraph.text
            budget.spend(text)
            blocks.append(ExtractedBlock(text=text, source_position={"paragraph": paragraph_number}))
        return blocks
    except KnowledgeError:
        raise
    except Exception:
        raise _parse_failed("DOCX 解析失败") from None


def _extract_plain_text(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    try:
        data = path.read_bytes()
    except OSError:
        raise _parse_failed("文本文件读取失败") from None
    text = _decode_text_bytes(data)
    budget.spend(text)
    return [ExtractedBlock(text=text, source_position={})]


def _extract_csv(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    try:
        data = path.read_bytes()
    except OSError:
        raise _parse_failed("CSV 文件读取失败") from None
    decoded = _decode_text_bytes(data)
    blocks: list[ExtractedBlock] = []
    try:
        for row_number, row in enumerate(csv.reader(io.StringIO(decoded)), start=1):
            cells = [cell.strip() for cell in row if cell.strip()]
            if cells:
                text = ", ".join(cells)
                budget.spend(text)
                blocks.append(ExtractedBlock(text=text, source_position={"row": row_number}))
    except csv.Error:
        raise _parse_failed("CSV 解析失败") from None
    return blocks


def _extract_xlsx(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    import openpyxl

    _ensure_zip_within_budget(path, budget, kind="XLSX")
    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        raise _parse_failed("XLSX 解析失败") from None
    try:
        blocks: list[ExtractedBlock] = []
        for sheet in workbook.worksheets:
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if cells:
                    text = ", ".join(cells)
                    budget.spend(text)
                    blocks.append(
                        ExtractedBlock(
                            text=text,
                            source_position={"sheet": sheet.title, "row": row_number},
                        )
                    )
        return blocks
    except KnowledgeError:
        raise
    except Exception:
        raise _parse_failed("XLSX 解析失败") from None
    finally:
        workbook.close()


_HTML_DROP_TAGS = ("script", "style", "noscript", "template")


def _html_text(markup: bytes | str) -> str:
    """Visible text of an HTML document via the stdlib parser.

    bs4 sniffs the encoding itself (meta charset/BOM) when given bytes, so
    HTML does not go through ``_decode_text_bytes``. Newlines separate the
    text nodes; the splitter's fallback separators handle the rest.
    """

    from bs4 import BeautifulSoup

    soup = BeautifulSoup(markup, "html.parser")
    for tag in soup(_HTML_DROP_TAGS):
        tag.decompose()
    lines = [line.strip() for line in soup.get_text("\n").splitlines()]
    return "\n".join(line for line in lines if line)


def _extract_html(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    try:
        data = path.read_bytes()
    except OSError:
        raise _parse_failed("HTML 文件读取失败") from None
    try:
        text = _html_text(data)
    except Exception:
        raise _parse_failed("HTML 解析失败") from None
    budget.spend(text)
    return [ExtractedBlock(text=text, source_position={})]


def _extract_pptx(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    from pptx import Presentation

    _ensure_zip_within_budget(path, budget, kind="PPTX")
    try:
        presentation = Presentation(str(path))
        blocks: list[ExtractedBlock] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            pieces: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    pieces.extend(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
                elif getattr(shape, "has_table", False) and shape.has_table:
                    for table_row in shape.table.rows:
                        cells = [cell.text.strip() for cell in table_row.cells if cell.text.strip()]
                        if cells:
                            pieces.append(", ".join(cells))
            text = "\n".join(pieces)
            budget.spend(text)
            blocks.append(ExtractedBlock(text=text, source_position={"slide": slide_number}))
        return blocks
    except KnowledgeError:
        raise
    except Exception:
        raise _parse_failed("PPTX 解析失败") from None


def _extract_epub(path: Path, budget: _CharBudget) -> list[ExtractedBlock]:
    import ebooklib
    from ebooklib import epub

    _ensure_zip_within_budget(path, budget, kind="EPUB")
    try:
        book = epub.read_epub(str(path), options={"ignore_ncx": True})
        blocks: list[ExtractedBlock] = []
        # Navigation documents are tables of contents, not book content.
        chapters = [item for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT) if item.is_chapter()]
        for chapter_number, item in enumerate(chapters, start=1):
            text = _html_text(item.get_content())
            budget.spend(text)
            blocks.append(ExtractedBlock(text=text, source_position={"chapter": chapter_number}))
        return blocks
    except KnowledgeError:
        raise
    except Exception:
        raise _parse_failed("EPUB 解析失败") from None


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".txt": _extract_plain_text,
    ".md": _extract_plain_text,
    ".csv": _extract_csv,
    ".xlsx": _extract_xlsx,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".pptx": _extract_pptx,
    ".epub": _extract_epub,
}
